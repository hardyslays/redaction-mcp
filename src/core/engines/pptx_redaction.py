"""python-pptx implementation of PPTX redaction."""

from __future__ import annotations

import base64
import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.core.models.document import Document
from src.core.models.redaction import (
    BoundingBox,
    BoundingBoxTarget,
    PageTarget,
    RedactionOptions,
    RedactionTarget,
    RegexTarget,
    TextTarget,
)
from src.core.services.text_matching import compile_text_pattern


def _text_frames(shapes: object) -> list[object]:
    frames = []
    for shape in shapes:  # type: ignore[union-attr]
        if shape.has_text_frame:
            frames.append(shape.text_frame)
        if shape.has_table:
            frames.extend(cell.text_frame for row in shape.table.rows for cell in row.cells)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            frames.extend(_text_frames(shape.shapes))
    return frames


def _shapes(shapes: object) -> list[object]:
    result = []
    for shape in shapes:  # type: ignore[union-attr]
        result.append(shape)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result.extend(_shapes(shape.shapes))
    return result


def _replacement(text: str, options: RedactionOptions) -> str:
    return "[REDACT]" if options.redaction_type == "mask" else "*" * len(text)


def _replace(frame: object, pattern: re.Pattern[str], options: RedactionOptions, limit: int = 0) -> int:
    text, replacements = pattern.subn(lambda match: _replacement(match.group(), options), frame.text, count=limit)
    if replacements:
        frame.text = text
    return replacements


def _slides(presentation: Presentation, selection: list[int] | None) -> list[object]:
    indexes = selection if selection is not None else range(len(presentation.slides))
    slides = []
    for index in indexes:
        if not 0 <= index < len(presentation.slides):
            raise ValueError(f"Slide {index} is out of range (presentation has {len(presentation.slides)} slides)")
        slides.append(presentation.slides[index])
    return slides


def _remove_shape(shape: object) -> None:
    element = shape._element  # type: ignore[attr-defined]
    element.getparent().remove(element)


def _intersects(shape: object, box: tuple[int, int, int, int]) -> bool:
    x, y, width, height = box
    left, top = shape.left, shape.top  # type: ignore[attr-defined]
    right, bottom = left + shape.width, top + shape.height  # type: ignore[attr-defined]
    return not (
        right <= x
        or x + width <= left
        or bottom <= y
        or y + height <= top
    )


def _rectangles_intersect(first: tuple[int, int, int, int], second: tuple[int, int, int, int]) -> bool:
    x, y, width, height = first
    other_x, other_y, other_width, other_height = second
    return not (
        x + width <= other_x
        or other_x + other_width <= x
        or y + height <= other_y
        or other_y + other_height <= y
    )


def _redact_words_in_box(slide: object, box: BoundingBox, presentation: Presentation, options: RedactionOptions) -> None:
    target = (
        int(box.x * presentation.slide_width),
        int(box.y * presentation.slide_height),
        int(box.width * presentation.slide_width),
        int(box.height * presentation.slide_height),
    )
    for shape in _shapes(slide.shapes):
        if not shape.has_text_frame or not _intersects(shape, target):
            continue
        text = shape.text_frame.text
        words = list(re.finditer(r"\S+", text))
        if not words:
            continue
        lines = text.splitlines() or [text]
        line_offsets: list[int] = []
        offset = 0
        for line in lines:
            line_offsets.append(offset)
            offset += len(line) + 1
        changed = False
        parts: list[str] = []
        cursor = 0
        for word in words:
            line_index = max(index for index, line_offset in enumerate(line_offsets) if line_offset <= word.start())
            line_start = line_offsets[line_index]
            line_length = len(lines[line_index])
            word_left = shape.left + shape.width * (word.start() - line_start) / line_length
            word_right = shape.left + shape.width * (word.end() - line_start) / line_length
            word_box = (
                int(word_left),
                int(shape.top + shape.height * line_index / len(lines)),
                int(word_right - word_left),
                int(shape.height / len(lines)),
            )
            if _rectangles_intersect(word_box, target):
                parts.extend((text[cursor:word.start()], _replacement(word.group(), options)))
                cursor = word.end()
                changed = True
        if changed:
            parts.append(text[cursor:])
            shape.text_frame.text = "".join(parts)


def _remove_all_shapes(slide: object) -> None:
    for shape in list(slide.shapes):
        _remove_shape(shape)


def redact_pptx_document(
    document: Document, targets: list[RedactionTarget], options: RedactionOptions, *, data: bytes | None = None
) -> Document:
    """Apply text redactions and redact words selected by area/page."""
    if data is None and document.base64data is None:
        raise ValueError("PPTX engine requires in-memory document base64data")
    presentation = Presentation(BytesIO(data if data is not None else document.decoded_bytes()))
    frames_by_slide = {index: _text_frames(slide.shapes) for index, slide in enumerate(presentation.slides)}
    for target in targets:
        if isinstance(target, TextTarget):
            pattern = compile_text_pattern(target.values, ignore_case=target.ignore_case, partial_match=target.partial_match)
            selection = target.pages if target.pages is not None else range(len(presentation.slides))
            for index in selection:
                if not 0 <= index < len(presentation.slides):
                    raise ValueError(f"Slide {index} is out of range (presentation has {len(presentation.slides)} slides)")
                for frame in frames_by_slide[index]:
                    _replace(frame, pattern, options)
        elif isinstance(target, RegexTarget):
            flags = re.IGNORECASE if target.ignore_case else 0
            if not target.allow_unicode:
                flags |= re.ASCII
            matched = False
            selection = target.pages if target.pages is not None else range(len(presentation.slides))
            for index in selection:
                if not 0 <= index < len(presentation.slides):
                    raise ValueError(f"Slide {index} is out of range (presentation has {len(presentation.slides)} slides)")
                frames = frames_by_slide[index]
                for pattern in (re.compile(value, flags) for value in target.patterns):
                    for frame in frames:
                        replaced = _replace(frame, pattern, options, limit=1 if target.only_first_match else 0)
                        if replaced and target.only_first_match:
                            matched = True
                            break
                    if matched:
                        break
        elif isinstance(target, BoundingBoxTarget):
            for box in target.values:
                slide = _slides(presentation, [box.page])[0]
                _redact_words_in_box(slide, box, presentation, options)
        elif isinstance(target, PageTarget):
            for slide in _slides(presentation, target.values):
                _remove_all_shapes(slide)
    output = BytesIO()
    presentation.save(output)
    filename = Path(document.filename or "document.pptx")
    return Document(
        base64data=base64.b64encode(output.getvalue()).decode("ascii"),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{filename.stem}-redacted.pptx",
    )
