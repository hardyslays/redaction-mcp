"""python-pptx implementation of permanent text data replacement."""

from __future__ import annotations

import base64
import re
from io import BytesIO
from pathlib import Path

from pptx import Presentation

from src.core.engines.pptx_redaction import _slides, _text_frames
from src.core.models.document import Document
from src.core.models.replacement import ReplacementTarget
from src.core.services.replacement_text import replacement_text


def _replace(frame: object, pattern: re.Pattern[str], target: ReplacementTarget) -> None:
    text, count = pattern.subn(lambda match: replacement_text(match.group(), target), frame.text)
    if count:
        # Assigning replacement text removes the source text from its XML runs.
        frame.text = text


def replace_pptx_document(document: Document, targets: list[ReplacementTarget], *, data: bytes | None = None) -> Document:
    """Replace exact text in text boxes, table cells, and grouped slide shapes."""
    if data is None and document.base64data is None:
        raise ValueError("PPTX engine requires in-memory document base64data")
    presentation = Presentation(BytesIO(data if data is not None else document.decoded_bytes()))
    frames_by_slide = {index: _text_frames(slide.shapes) for index, slide in enumerate(presentation.slides)}
    for target in targets:
        flags = re.IGNORECASE if target.ignore_case else 0
        pattern = re.compile("|".join(re.escape(value) for value in sorted(target.values, key=len, reverse=True)), flags)
        selection = target.pages if target.pages is not None else range(len(presentation.slides))
        for index in selection:
            if not 0 <= index < len(presentation.slides):
                raise ValueError(f"Slide {index} is out of range (presentation has {len(presentation.slides)} slides)")
            for frame in frames_by_slide[index]:
                _replace(frame, pattern, target)
    output = BytesIO()
    presentation.save(output)
    filename = Path(document.filename or "document.pptx")
    return Document(
        base64data=base64.b64encode(output.getvalue()).decode("ascii"),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{filename.stem}-replaced.pptx",
    )
