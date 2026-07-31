"""Bounded document-text extraction for MCP agents."""

from __future__ import annotations

from io import BytesIO

import fitz
from docx import Document as DocxDocument
from pydantic import BaseModel, Field
from pptx import Presentation

from src.core.models.document import Document
from src.core.services.mime_detection import (
    DOCX_MIME_TYPE,
    PDF_MIME_TYPE,
    PPTM_MIME_TYPE,
    PPTX_MIME_TYPE,
    TEXT_MIME_TYPE,
    detect_mime_type,
)


class DocumentPage(BaseModel):
    """Text visible to an agent on one zero-based document page or slide."""

    index: int = Field(ge=0)
    text: str


class DocumentInspection(BaseModel):
    """Bounded text extraction that lets agents select redaction targets."""

    filename: str
    mime_type: str
    pages: list[DocumentPage]
    truncated: bool


def _pptx_text(data: bytes) -> list[str]:
    presentation = Presentation(BytesIO(data))
    return ["\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)) for slide in presentation.slides]


def inspect_document(document: Document, max_characters: int) -> DocumentInspection:
    """Extract visible text without returning more than ``max_characters``."""
    data = document.decoded_bytes()
    mime_type = document.mime_type or detect_mime_type(data, document.filename)
    if mime_type == PDF_MIME_TYPE:
        pdf = fitz.open(stream=data, filetype="pdf")
        try:
            texts = [page.get_text() for page in pdf]
        finally:
            pdf.close()
    elif mime_type == DOCX_MIME_TYPE:
        texts = ["\n".join(paragraph.text for paragraph in DocxDocument(BytesIO(data)).paragraphs)]
    elif mime_type in (PPTX_MIME_TYPE, PPTM_MIME_TYPE):
        texts = _pptx_text(data)
    elif mime_type == TEXT_MIME_TYPE:
        texts = [data.decode("utf-8")]
    else:
        raise ValueError("Only PDF, DOCX, PPTX, PPTM, and TXT documents are currently supported")

    remaining = max_characters
    pages: list[DocumentPage] = []
    truncated = False
    for index, text in enumerate(texts):
        if remaining <= 0:
            truncated = True
            break
        if len(text) > remaining:
            pages.append(DocumentPage(index=index, text=text[:remaining]))
            truncated = True
            break
        pages.append(DocumentPage(index=index, text=text))
        remaining -= len(text)
    return DocumentInspection(
        filename=document.filename or "document",
        mime_type=mime_type or "application/octet-stream",
        pages=pages,
        truncated=truncated,
    )
