import base64
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import fitz
import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from src.core.models.document import Document
from src.core.models.errors import RedactionError
from src.core.models.redaction import BoundingBox, BoundingBoxTarget, PageTarget, RedactionOptions, RegexTarget, TextTarget
from src.core.services.redaction_service import redact_document


def pdf_with_text(text: str = "secret visible") -> bytes:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), text)
    data = pdf.tobytes()
    pdf.close()
    return data


def docx_with_text(text: str = "secret visible") -> bytes:
    document = DocxDocument()
    document.add_paragraph(text)
    data = BytesIO()
    document.save(data)
    return data.getvalue()


def pptx_with_text(text: str = "secret visible") -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = text
    data = BytesIO()
    presentation.save(data)
    return data.getvalue()


def pptm_with_text(text: str) -> bytes:
    source = BytesIO(pptx_with_text(text))
    output = BytesIO()
    with ZipFile(source) as archive, ZipFile(output, "w", ZIP_DEFLATED) as replacement:
        for entry in archive.infolist():
            data = archive.read(entry.filename)
            if entry.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    b"application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
                )
            replacement.writestr(entry, data)
    return output.getvalue()


def output_text(result: Document, extension: str) -> str:
    data = base64.b64decode(result.base64data)
    if extension == "pdf":
        output = fitz.open(stream=data, filetype="pdf")
        try:
            return output[0].get_text()
        finally:
            output.close()
    if extension == "docx":
        return DocxDocument(BytesIO(data)).paragraphs[0].text
    return Presentation(BytesIO(data)).slides[0].shapes[0].text


def test_text_redaction_removes_matching_text() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pdf_with_text()).decode(), filename="source.pdf"),
        [TextTarget(type="text", values=["secret"])],
    )

    assert isinstance(result, Document)
    assert result.filename == "source-redacted.pdf"
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    assert "secret" not in output[0].get_text()
    assert "visible" in output[0].get_text()
    output.close()


def test_text_redaction_ignore_case_is_opt_in_for_pdf() -> None:
    source = Document(base64data=base64.b64encode(pdf_with_text("SECRET visible")).decode(), filename="source.pdf")

    case_sensitive = redact_document(source, [TextTarget(type="text", values=["secret"])])
    assert isinstance(case_sensitive, Document)
    output = fitz.open(stream=base64.b64decode(case_sensitive.base64data), filetype="pdf")
    assert "SECRET" in output[0].get_text()
    output.close()

    result = redact_document(source, [TextTarget(type="text", values=["secret"], ignore_case=True)])
    assert isinstance(result, Document)
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    assert "SECRET" not in output[0].get_text()
    output.close()


@pytest.mark.parametrize("extension", ["pdf", "docx", "pptx"])
def test_text_redaction_requires_alphanumeric_boundaries_by_default(extension: str) -> None:
    source = "john johnathon"
    builders = {"pdf": pdf_with_text, "docx": docx_with_text, "pptx": pptx_with_text}
    document = Document(base64data=base64.b64encode(builders[extension](source)).decode(), filename=f"source.{extension}")

    safe_result = redact_document(document, [TextTarget(type="text", values=["john"])])
    assert isinstance(safe_result, Document)
    safe_text = output_text(safe_result, extension)
    assert "johnathon" in safe_text
    assert "john " not in safe_text

    partial_result = redact_document(document, [TextTarget(type="text", values=["john"], partial_match=True)])
    assert isinstance(partial_result, Document)
    assert "johnathon" not in output_text(partial_result, extension)


def test_normalized_bounding_box_redacts_area() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pdf_with_text("hide")).decode(), filename="source.pdf"),
        [BoundingBoxTarget(type="bounding_box", values=[BoundingBox(page=0, x=0, y=0, width=1, height=1)])],
    )

    assert isinstance(result, Document)
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    assert output[0].get_text().strip() == ""
    output.close()


def test_unsupported_document_returns_typed_error() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(b"plain text").decode(), filename="note.csv"),
        [TextTarget(type="text", values=["text"])],
    )

    assert isinstance(result, RedactionError)
    assert result.message == "Only PDF, DOCX, PPTX, PPTM, and TXT documents are currently supported"


def test_txt_text_redaction_replaces_matching_text_with_asterisks() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(b"secret visible").decode(), filename="source.txt"),
        [TextTarget(type="text", values=["secret"])],
    )

    assert isinstance(result, Document)
    assert result.filename == "source-redacted.txt"
    assert result.mime_type == "text/plain"
    assert base64.b64decode(result.base64data) == b"****** visible"


def test_txt_page_redaction_masks_text_and_preserves_whitespace_layout() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(b"Name: Jane Doe\n\tID 42\n").decode(), filename="source.txt"),
        [PageTarget(type="page", values=[0])],
    )

    assert isinstance(result, Document)
    assert base64.b64decode(result.base64data) == b"***** **** ***\n\t** **\n"


@pytest.mark.parametrize(
    ("target", "message"),
    [
        (PageTarget(type="page", values=[1]), "TXT redaction supports only page index 0"),
        (
            BoundingBoxTarget(type="bounding_box", values=[BoundingBox(page=0, x=0, y=0, width=1, height=1)]),
            "TXT redaction does not support bounding-box targets",
        ),
        (TextTarget(type="text", values=["secret"], pages=[1]), "TXT redaction supports only page index 0"),
    ],
)
def test_txt_redaction_rejects_unsupported_targets_or_pages(target: object, message: str) -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(b"secret").decode(), filename="source.txt"), [target]
    )

    assert isinstance(result, RedactionError)
    assert result.message == message


def test_docx_text_redaction_removes_matching_text() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text()).decode(), filename="source.docx"),
        [TextTarget(type="text", values=["secret"])],
    )

    assert isinstance(result, Document)
    assert result.filename == "source-redacted.docx"
    output = DocxDocument(BytesIO(base64.b64decode(result.base64data)))
    assert "secret" not in output.paragraphs[0].text
    assert "******" in output.paragraphs[0].text
    assert "visible" in output.paragraphs[0].text


def test_docx_page_redaction_returns_typed_error() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text()).decode(), filename="source.docx"),
        [PageTarget(type="page", values=[0])],
    )

    assert isinstance(result, RedactionError)
    assert result.message == "DOCX redaction supports only text and regex targets"


def test_docx_regex_redaction_removes_matching_text() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text("ID AB123456 visible")).decode(), filename="source.docx"),
        [RegexTarget(type="regex", patterns=[r"[A-Z]{2}\d{6}"], ignore_case=False)],
    )

    assert isinstance(result, Document)
    output = DocxDocument(BytesIO(base64.b64decode(result.base64data)))
    assert "AB123456" not in output.paragraphs[0].text
    assert "visible" in output.paragraphs[0].text


def test_pptx_text_redaction_removes_matching_text() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pptx_with_text()).decode(), filename="source.pptx"),
        [TextTarget(type="text", values=["secret"])],
    )

    assert isinstance(result, Document)
    assert result.filename == "source-redacted.pptx"
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert "secret" not in output.slides[0].shapes[0].text
    assert "******" in output.slides[0].shapes[0].text
    assert "visible" in output.slides[0].shapes[0].text


def test_pptm_text_redaction_removes_matching_text() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pptm_with_text("secret visible")).decode(), filename="source.pptm"),
        [TextTarget(type="text", values=["secret"])],
    )

    assert isinstance(result, Document)
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert "secret" not in output.slides[0].shapes[0].text


def test_docx_and_pptx_text_redaction_ignore_case() -> None:
    docx_result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text("SECRET visible")).decode(), filename="source.docx"),
        [TextTarget(type="text", values=["secret"], ignore_case=True)],
    )
    assert isinstance(docx_result, Document)
    docx = DocxDocument(BytesIO(base64.b64decode(docx_result.base64data)))
    assert "SECRET" not in docx.paragraphs[0].text

    pptx_result = redact_document(
        Document(base64data=base64.b64encode(pptx_with_text("SECRET visible")).decode(), filename="source.pptx"),
        [TextTarget(type="text", values=["secret"], ignore_case=True)],
    )
    assert isinstance(pptx_result, Document)
    pptx = Presentation(BytesIO(base64.b64decode(pptx_result.base64data)))
    assert "SECRET" not in pptx.slides[0].shapes[0].text


def test_mask_redaction_replaces_text_with_redact_marker() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text()).decode(), filename="source.docx"),
        [TextTarget(type="text", values=["secret"])],
        RedactionOptions(redaction_type="mask"),
    )

    assert isinstance(result, Document)
    output = DocxDocument(BytesIO(base64.b64decode(result.base64data)))
    assert "secret" not in output.paragraphs[0].text
    assert "[REDACT]" in output.paragraphs[0].text


def test_docx_bounding_box_redacts_only_intersecting_words() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(docx_with_text("secret visible")).decode(), filename="source.docx"),
        [BoundingBoxTarget(type="bounding_box", values=[BoundingBox(page=0, x=0, y=0, width=0.5, height=1)])],
    )

    assert isinstance(result, Document)
    assert output_text(result, "docx") == "****** visible"


def test_pptx_bounding_box_redacts_only_intersecting_words() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pptx_with_text("secret visible")).decode(), filename="source.pptx"),
        [BoundingBoxTarget(type="bounding_box", values=[BoundingBox(page=0, x=0.07, y=0.07, width=0.16, height=0.14)])],
    )

    assert isinstance(result, Document)
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert len(output.slides[0].shapes) == 1
    assert output.slides[0].shapes[0].text == "****** visible"


@pytest.mark.parametrize(
    "values",
    [
        {"page": 0, "x": 0, "y": 0, "width": 1, "height": 1, "units": "pixels"},
        {"page": 0, "x": 0.8, "y": 0, "width": 0.3, "height": 1},
    ],
)
def test_bounding_box_rejects_non_normalized_or_out_of_page_values(values: dict[str, object]) -> None:
    with pytest.raises(ValueError):
        BoundingBox.model_validate(values)


def test_pptx_page_redaction_removes_all_slide_shapes() -> None:
    result = redact_document(
        Document(base64data=base64.b64encode(pptx_with_text()).decode(), filename="source.pptx"),
        [PageTarget(type="page", values=[0])],
    )

    assert isinstance(result, Document)
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert len(output.slides[0].shapes) == 0


def test_pdf_regex_redacts_exact_multiline_span_without_redacting_surrounding_text() -> None:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Owner: Jane")
    page.insert_text((72, 96), "Doe; public")
    data = pdf.tobytes()
    pdf.close()

    result = redact_document(
        Document(base64data=base64.b64encode(data).decode(), filename="source.pdf"),
        [RegexTarget(type="regex", patterns=[r"Jane\s+Doe"], ignore_case=False)],
    )

    assert isinstance(result, Document)
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    text = output[0].get_text()
    assert "Jane" not in text and "Doe" not in text
    assert "Owner:" in text and "public" in text
    output.close()


def test_regex_only_first_match_is_document_wide_for_pdf_and_pptx() -> None:
    pdf = fitz.open()
    first = pdf.new_page()
    first.insert_text((72, 72), "ID AB123456")
    second = pdf.new_page()
    second.insert_text((72, 72), "ID CD234567")
    data = pdf.tobytes()
    pdf.close()
    target = RegexTarget(type="regex", patterns=[r"[A-Z]{2}\d{6}"], ignore_case=False, only_first_match=True)

    result = redact_document(Document(base64data=base64.b64encode(data).decode(), filename="source.pdf"), [target])
    assert isinstance(result, Document)
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    assert "AB123456" not in output[0].get_text()
    assert "CD234567" in output[1].get_text()
    output.close()

    pptx_result = redact_document(
        Document(base64data=base64.b64encode(pptx_with_text("AB123456 AB123456")).decode(), filename="source.pptx"),
        [target],
    )
    assert isinstance(pptx_result, Document)
    presentation = Presentation(BytesIO(base64.b64decode(pptx_result.base64data)))
    assert presentation.slides[0].shapes[0].text.count("AB123456") == 1
