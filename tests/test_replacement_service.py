import base64
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import fitz
import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches

from src.core.engines.pdf_replacement import _replacement_lines, replacement_text
from src.core.models.document import Document
from src.core.models.errors import ReplacementError
from src.core.models.replacement import TextReplacementTarget
from src.core.services.replacement_service import replace_document


def pdf_with_text(text: str) -> bytes:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), text)
    data = pdf.tobytes()
    pdf.close()
    return data


def target(kind: str, source: str = "Jane Doe", **values: object) -> TextReplacementTarget:
    return TextReplacementTarget(type="text", values=[source], replacement_type=kind, **values)


def docx_with_text(text: str) -> bytes:
    document = DocxDocument()
    document.add_paragraph(text)
    data = BytesIO()
    document.save(data)
    return data.getvalue()


def pptx_with_text(text: str) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    shape.text = text
    data = BytesIO()
    presentation.save(data)
    return data.getvalue()


def pptm_with_text(text: str) -> bytes:
    source = BytesIO(pptx_with_text(text))
    output = BytesIO()
    with ZipFile(source) as archive, ZipFile(output, "w", ZIP_DEFLATED) as replacement:
        for entry in archive.infolist():
            data = archive.read(entry.filename)
            if entry.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                    b"application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml",
                )
            replacement.writestr(entry, data)
    return output.getvalue()


def output_text(result: Document, extension: str) -> str:
    data = base64.b64decode(result.base64data)
    if extension == "pdf":
        output = fitz.open(stream=data, filetype="pdf")
        try:
            return output[0].get_text()
        finally:
            output.close()
    if extension == "docx":
        return DocxDocument(BytesIO(data)).paragraphs[0].text
    return Presentation(BytesIO(data)).slides[0].shapes[0].text


def test_partial_replacement_masks_non_whitespace_and_keeps_suffix() -> None:
    assert replacement_text("Jane Doe", target("PARTIAL")) == "**** *oe"
    assert replacement_text("AB CD", target("PARTIAL")) == "** CD"
    assert replacement_text("Jane\nDoe", target("PARTIAL")) == "****\n*oe"


def test_static_replacement_is_bracketed() -> None:
    assert replacement_text("Jane Doe", target("STATIC", static_text="REDACTED")) == "[REDACTED]"
    assert _replacement_lines("Long\nline", "[SAFE]", static=True) == ["[SA", "FE]"]


def test_regex_replacement_preserves_character_pattern_without_source_text() -> None:
    assert replacement_text("AB-129 z", target("REGEX")) == "BC-230 a"


def test_pdf_replacement_permanently_removes_original_and_inserts_static_value() -> None:
    result = replace_document(
        Document(base64data=base64.b64encode(pdf_with_text("Owner Jane Doe")).decode(), filename="source.pdf"),
        [target("STATIC", static_text="REDACTED")],
    )

    assert not isinstance(result, ReplacementError)
    assert result.filename == "source-replaced.pdf"
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    text = output[0].get_text()
    assert "Jane Doe" not in text
    assert "[REDACTED]" in text
    output.close()


def test_pdf_replacement_ignore_case_is_opt_in() -> None:
    source = Document(base64data=base64.b64encode(pdf_with_text("Owner JANE DOE")).decode(), filename="source.pdf")

    case_sensitive = replace_document(source, [target("STATIC", static_text="SAFE")])
    assert not isinstance(case_sensitive, ReplacementError)
    output = fitz.open(stream=base64.b64decode(case_sensitive.base64data), filetype="pdf")
    assert "JANE DOE" in output[0].get_text()
    output.close()

    result = replace_document(source, [target("STATIC", static_text="SAFE", ignore_case=True)])
    assert not isinstance(result, ReplacementError)
    output = fitz.open(stream=base64.b64decode(result.base64data), filetype="pdf")
    assert "JANE DOE" not in output[0].get_text()
    assert "[SAFE]" in output[0].get_text()
    output.close()


@pytest.mark.parametrize("extension", ["pdf", "docx", "pptx"])
def test_replacement_requires_alphanumeric_boundaries_by_default(extension: str) -> None:
    source = "john johnathon"
    builders = {"pdf": pdf_with_text, "docx": docx_with_text, "pptx": pptx_with_text}
    document = Document(base64data=base64.b64encode(builders[extension](source)).decode(), filename=f"source.{extension}")

    safe_result = replace_document(document, [target("STATIC", source="john", static_text="SAFE")])
    assert not isinstance(safe_result, ReplacementError)
    safe_text = output_text(safe_result, extension)
    assert "johnathon" in safe_text
    assert safe_text.count("[SAFE]") == 1

    partial_result = replace_document(
        document, [target("STATIC", source="john", static_text="SAFE", partial_match=True)]
    )
    assert not isinstance(partial_result, ReplacementError)
    assert "johnathon" not in output_text(partial_result, extension)


def test_docx_replacement_permanently_removes_original_for_all_strategies() -> None:
    for kind, expected, kwargs in (
        ("PARTIAL", "**** *oe", {}),
        ("STATIC", "[SAFE]", {"static_text": "SAFE"}),
        ("REGEX", "Kbof Epf", {}),
    ):
        result = replace_document(
            Document(base64data=base64.b64encode(docx_with_text("Owner Jane Doe")).decode(), filename="source.docx"),
            [target(kind, **kwargs)],
        )

        assert not isinstance(result, ReplacementError)
        assert result.filename == "source-replaced.docx"
        output = DocxDocument(BytesIO(base64.b64decode(result.base64data)))
        assert "Jane Doe" not in output.paragraphs[0].text
        assert expected in output.paragraphs[0].text


def test_docx_replacement_handles_multiline_text_and_rejects_page_filter() -> None:
    result = replace_document(
        Document(base64data=base64.b64encode(docx_with_text("Jane\nDoe")).decode(), filename="source.docx"),
        [target("PARTIAL", source="Jane\nDoe")],
    )

    assert not isinstance(result, ReplacementError)
    output = DocxDocument(BytesIO(base64.b64decode(result.base64data)))
    assert output.paragraphs[0].text == "****\n*oe"

    filtered = replace_document(
        Document(base64data=base64.b64encode(docx_with_text("Jane Doe")).decode(), filename="source.docx"),
        [TextReplacementTarget(type="text", values=["Jane Doe"], replacement_type="PARTIAL", pages=[0])],
    )
    assert isinstance(filtered, ReplacementError)
    assert filtered.message == "DOCX data replacement does not support page-restricted targets"


def test_pptx_replacement_permanently_removes_original_for_all_strategies() -> None:
    for kind, expected, kwargs in (
        ("PARTIAL", "**** *oe", {}),
        ("STATIC", "[SAFE]", {"static_text": "SAFE"}),
        ("REGEX", "Kbof Epf", {}),
    ):
        result = replace_document(
            Document(base64data=base64.b64encode(pptx_with_text("Owner Jane Doe")).decode(), filename="source.pptx"),
            [target(kind, **kwargs)],
        )

        assert not isinstance(result, ReplacementError)
        assert result.filename == "source-replaced.pptx"
        output = Presentation(BytesIO(base64.b64decode(result.base64data)))
        text = output.slides[0].shapes[0].text
        assert "Jane Doe" not in text
        assert expected in text


def test_pptm_replacement_removes_matching_text() -> None:
    result = replace_document(
        Document(base64data=base64.b64encode(pptm_with_text("Owner Jane Doe")).decode(), filename="source.pptm"),
        [target("STATIC", static_text="SAFE")],
    )

    assert not isinstance(result, ReplacementError)
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert "Jane Doe" not in output.slides[0].shapes[0].text
    assert "[SAFE]" in output.slides[0].shapes[0].text


def test_docx_and_pptx_replacement_ignore_case() -> None:
    docx_result = replace_document(
        Document(base64data=base64.b64encode(docx_with_text("Owner JANE DOE")).decode(), filename="source.docx"),
        [target("STATIC", static_text="SAFE", ignore_case=True)],
    )
    assert not isinstance(docx_result, ReplacementError)
    docx = DocxDocument(BytesIO(base64.b64decode(docx_result.base64data)))
    assert "JANE DOE" not in docx.paragraphs[0].text
    assert "[SAFE]" in docx.paragraphs[0].text

    pptx_result = replace_document(
        Document(base64data=base64.b64encode(pptx_with_text("Owner JANE DOE")).decode(), filename="source.pptx"),
        [target("STATIC", static_text="SAFE", ignore_case=True)],
    )
    assert not isinstance(pptx_result, ReplacementError)
    pptx = Presentation(BytesIO(base64.b64decode(pptx_result.base64data)))
    assert "JANE DOE" not in pptx.slides[0].shapes[0].text
    assert "[SAFE]" in pptx.slides[0].shapes[0].text


def test_pptx_replacement_handles_multiline_text_and_slide_filter() -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Jane\nDoe"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = "Jane Doe"
    data = BytesIO()
    presentation.save(data)

    result = replace_document(
        Document(base64data=base64.b64encode(data.getvalue()).decode(), filename="source.pptx"),
        [TextReplacementTarget(type="text", values=["Jane\nDoe"], replacement_type="PARTIAL", pages=[0])],
    )

    assert not isinstance(result, ReplacementError)
    output = Presentation(BytesIO(base64.b64decode(result.base64data)))
    assert output.slides[0].shapes[0].text == "****\n*oe"
    assert output.slides[1].shapes[0].text == "Jane Doe"


def test_replacement_rejects_unsupported_document() -> None:
    result = replace_document(
        Document(base64data=base64.b64encode(b"text").decode(), filename="note.csv"),
        [target("PARTIAL")],
    )

    assert isinstance(result, ReplacementError)
    assert result.message == "Only PDF, DOCX, PPTX, PPTM, and TXT documents are currently supported for data replacement"


@pytest.mark.parametrize(
    ("kind", "expected", "kwargs"),
    [
        ("PARTIAL", "**** *oe", {}),
        ("STATIC", "[SAFE]", {"static_text": "SAFE"}),
        ("REGEX", "Kbof Epf", {}),
    ],
)
def test_txt_replacement_supports_all_strategies(kind: str, expected: str, kwargs: dict[str, str]) -> None:
    result = replace_document(
        Document(base64data=base64.b64encode(b"Owner Jane Doe").decode(), filename="source.txt"),
        [target(kind, **kwargs)],
    )

    assert isinstance(result, Document)
    assert result.filename == "source-replaced.txt"
    assert result.mime_type == "text/plain"
    output = base64.b64decode(result.base64data).decode()
    assert "Jane Doe" not in output
    assert expected in output


def test_txt_replacement_accepts_page_zero_and_rejects_other_pages() -> None:
    source = Document(base64data=base64.b64encode(b"Jane Doe").decode(), filename="source.txt")

    valid = replace_document(source, [TextReplacementTarget(type="text", values=["Jane Doe"], replacement_type="PARTIAL", pages=[0])])
    assert isinstance(valid, Document)

    invalid = replace_document(source, [TextReplacementTarget(type="text", values=["Jane Doe"], replacement_type="PARTIAL", pages=[1])])
    assert isinstance(invalid, ReplacementError)
    assert invalid.message == "TXT data replacement supports only page index 0"
